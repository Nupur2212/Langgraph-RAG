import os
from gtts import gTTS
from pydub import AudioSegment
import fitz
from langchain.schema import Document
import chromadb
import chromadb.utils.embedding_functions as embedding_functions
from pptx import Presentation
from bs4 import BeautifulSoup
import google.generativeai as genai
from typing import TypedDict, List
from langchain_core.messages import BaseMessage, SystemMessage, HumanMessage, AIMessage
from langchain_core.prompts import ChatPromptTemplate
from pydantic import BaseModel, Field
from dotenv import load_dotenv

load_dotenv()
apikey=os.getenv("GOOGLE_API_KEY")

def load_files(folder_path):
    print(os.listdir(folder_path))
    documents=[]
    for filename in os.listdir(folder_path):
        print(filename)
        file_path = os.path.join(folder_path, filename)
        try:
            if filename.endswith(".pdf"):
                doc = fitz.open(file_path)
                text = ""
                for page in doc:
                    text += page.get_text()
                documents.append(Document(page_content=text, metadata={"source": file_path}))
            elif filename.endswith(".txt"):
                with open(file_path, 'r', encoding='utf-8') as file:
                    text = file.read()
                    documents.append(Document(page_content=text, metadata={"source": file_path}))
            elif filename.endswith('.pptx'):
                presentation = Presentation(file_path)
                text = []
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text.append(shape.text)
                text='\n'.join(text)
                documents.append(Document(page_content=text, metadata={"source": file_path}))
            elif filename.endswith('.html'):
                with open(file_path, 'r', encoding='utf-8') as f:
                    html_content = f.read()
                soup = BeautifulSoup(html_content, 'html.parser')
                text = soup.get_text()
                documents.append(Document(page_content=text, metadata={"source": file_path}))
        except Exception as e:
            print("Error loading :",filename,e)
    return documents

def split_text_recursively(text, max_length=500, chunk_overlap=20):
    chunks = []
    start = 0
    while start < len(text):
        end = start + max_length
        if end < len(text):
            end = text.rfind(' ', start, end) + 1
            if end <= start:
                end = start + max_length
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        start = end - chunk_overlap
        if start >= len(text):
            break
    return chunks

def chunks(extracted_data):
    text_chunks = []
    for doc in extracted_data:
        chunks = split_text_recursively(doc.page_content, max_length=100, chunk_overlap=20)
        for chunk in chunks:
            text_chunks.append(Document(page_content=chunk, metadata=doc.metadata))
    client = chromadb.PersistentClient(path='embeddings/gemini')
    # client.delete_collection('pdf_rag')
    google_ef = embedding_functions.GoogleGenerativeAiEmbeddingFunction(api_key=apikey)
    collection = client.get_or_create_collection(name='pdf_rag', embedding_function=google_ef)
    collection.add(documents=[d.page_content for d in text_chunks], 
                    metadatas=[d.metadata for d in text_chunks], 
                    ids=[str(i) for i in range(len(text_chunks))])
    return google_ef,collection
    
def build_escaped_context(context):
    escaped_context = ''
    sources = []
    for item in context:
        escaped_context += item['document'] + "\n"
        sources.append(item['metadata'])
    return escaped_context, sources

def find_relevant_context(query, db, n_result=3):
    results = db.query(query_texts=[query], n_results=n_result)
    documents = results['documents'][0]
    metadatas = results['metadatas'][0]
    context = [{'document': doc, 'metadata': meta} for doc, meta in zip(documents, metadatas)]
    escaped_context, source = build_escaped_context(context)
    return escaped_context, source

def create_prompt_for_gemini(query, context,history):
    prompt = f"""
    You are an agent that answers questions using the text from the context below.
    Both the question and context is shared with you and you should answer the question on the basis of the context and not hallucinate.
    If the context does not have enough information for you, inform about the absence of relevant context as part of your answer.

    Context : {context}

    Question : {query}

    History: {history} <messages of the converstaion with you>

    Answer :
    """
    return prompt


def generate_answer_from_gemini(prompt):
    model = genai.GenerativeModel('gemini-1.5-pro-latest')
    result = model.generate_content(prompt)
    return result

documents=load_files(folder_path=os.getcwd()+'/files')
print(len(documents))
google_ef,collection=chunks(documents)
# query = input()
# results, sources = find_relevant_context(query, collection, 3)
# prompt, sources = create_prompt_for_gemini(query, results, sources)
# answer_text = generate_answer_from_gemini(prompt)

class AgentState(TypedDict):
    messages: List[BaseMessage]
    documents: List[Document]=[]
    # sources:List[str]
    on_topic: str
    rephrased_question: str
    proceed_to_generate: bool
    rephrase_count: int=0
    question: HumanMessage
    # db: 

def retrievenode(state:AgentState)->AgentState:
    global collection
    if "messages" not in state or state["messages"] is None:
        state['rephrase_count']=0
        state["messages"] = []

    if state["question"] not in state["messages"]:
        state["messages"].append(state["question"])

    if len(state["messages"]) > 1:
        conversation = state["messages"][:-1]
        current_question = state["question"].content
        messages = [
            SystemMessage(
                content="You are a helpful assistant that rephrases the user's question to be a standalone question optimized for retrieval."
            )
        ]
        messages.extend(conversation)
        messages.append(HumanMessage(content=current_question))

    query=state['question'].content
    results, sources = find_relevant_context(query, collection, 3)
    if "documents" not in state or state["documents"] is []:
        print("in")
        state['documents']=[]
    print(state)
    for r,s in zip(results,sources):
        state['documents'].append(Document(page_content=r,metadata=s))
    return state

class GradeDocument(BaseModel):
    score: str = Field(
        description="Document is relevant to the question? If yes -> 'Yes' if not -> 'No'"
    )

def retrieval_grader(state: AgentState):
    print("Entering retrieval_grader")
    system_message = SystemMessage(
        content="""You are a grader assessing the relevance of a retrieved document to a user question.
        Only answer with 'Yes' or 'No'.

        If the document contains information relevant to the user's question, respond with 'Yes'.
        Otherwise, respond with 'No'."""
        )

    llm = genai.GenerativeModel('gemini-1.5-pro-latest')
    structured_llm = llm.with_structured_output(GradeDocument)

    relevant_docs = []
    for doc in state["documents"]:
        human_message = HumanMessage(
            content=f"User question: {state['question']}\n\nRetrieved document:\n{doc.page_content}"
        )
        grade_prompt = ChatPromptTemplate.from_messages([system_message, human_message])
        grader_llm = grade_prompt | structured_llm
        result = grader_llm.invoke({})
        print(
            f"Grading document: {doc.metadata}... Result: {result.score.strip()}"
        )
        if result.score.strip().lower() == "yes":
            relevant_docs.append(doc)
    state["documents"] = relevant_docs
    state["proceed_to_generate"] = len(relevant_docs) > 0
    print(f"retrieval_grader: proceed_to_generate = {state['proceed_to_generate']}")
    return state

def proceed_router(state: AgentState):
    print("Entering proceed_router")
    rephrase_count = state.get("rephrase_count", 0)
    if state.get("proceed_to_generate", False):
        print("Routing to generate_answer")
        return "generate answer"
    elif rephrase_count >= 2:
        print("Maximum rephrase attempts reached. Cannot find relevant documents.")
        return "cannot answer"
    else:
        state['rephrase_count']+=1
        return "continue"
    
def generatenode(state:AgentState)->AgentState:
    print("Entering generate_answer")
    if "messages" not in state or state["messages"] is None:
        raise ValueError("State must include 'messages' before generating an answer.")

    history = state["messages"]
    documents = state["documents"]
    rephrased_question = state["question"]

    prompt= create_prompt_for_gemini(rephrased_question, state['documents'],history)
    answer_text = generate_answer_from_gemini(prompt)

    # response = rag_chain.invoke(
    #     {"history": history, "context": documents, "question": rephrased_question}
    # )

    # generation = response.content.strip()

    state["messages"].append(AIMessage(content=answer_text))
    print(f"generate_answer: Generated response: {answer_text}")
    return state

from langgraph.graph import StateGraph, END

graph=StateGraph(AgentState)

graph.add_node('retrieve',retrievenode)
graph.add_node('generate',generatenode)

graph.set_entry_point('retrieve')

graph.add_conditional_edges('retrieve',proceed_router,{"generate answer":'generate',"cannot answer":END,"continue":'retrieve'})

graph.add_edge('retrieve','generate')
graph.add_edge('generate',END)

app=graph.compile()

input_data = {"question": HumanMessage(content="How is Diwali celebrated?")}
print(app.invoke(input=input_data, config={"configurable": {"thread_id": 1}}))